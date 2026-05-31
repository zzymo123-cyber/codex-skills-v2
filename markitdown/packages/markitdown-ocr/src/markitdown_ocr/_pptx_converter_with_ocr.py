"""
Enhanced PPTX Converter with improved OCR support.
Already has LLM-based image description, this enhances it with traditional OCR fallback.
"""

import io
import sys
from typing import Any, BinaryIO, Optional

from typing import BinaryIO, Any, Optional

from markitdown.converters import HtmlConverter
from markitdown import DocumentConverter, DocumentConverterResult, StreamInfo
from markitdown._exceptions import (
    MissingDependencyException,
    MISSING_DEPENDENCY_MESSAGE,
)
from ._ocr_service import LLMVisionOCRService

_dependency_exc_info = None
try:
    import pptx
except ImportError:
    _dependency_exc_info = sys.exc_info()


class PptxConverterWithOCR(DocumentConverter):
    """Enhanced PPTX Converter with OCR fallback."""

    def __init__(self, ocr_service: Optional[LLMVisionOCRService] = None):
        super().__init__()
        self._html_converter = HtmlConverter()
        self.ocr_service = ocr_service

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension == ".pptx":
            return True

        if mimetype.startswith(
            "application/vnd.openxmlformats-officedocument.presentationml"
        ):
            return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> DocumentConverterResult:
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[1].with_traceback(
                _dependency_exc_info[2]
            )  # type: ignore[union-attr]

        # Get OCR service (from kwargs or instance)
        ocr_service: Optional[LLMVisionOCRService] = (
            kwargs.get("ocr_service") or self.ocr_service
        )
        llm_client = kwargs.get("llm_client")

        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += f"\\n\\n<!-- Slide number: {slide_num} -->\\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content

                # Pictures
                if self._is_picture(shape):
                    # Get image data
                    image_stream = io.BytesIO(shape.image.blob)

                    # Try LLM description first if available
                    llm_description = ""
                    if llm_client and kwargs.get("llm_model"):
                        try:
                            from ._llm_caption import llm_caption

                            image_filename = shape.image.filename
                            image_extension = None
                            if image_filename:
                                import os

                                image_extension = os.path.splitext(image_filename)[1]

                            image_stream_info = StreamInfo(
                                mimetype=shape.image.content_type,
                                extension=image_extension,
                                filename=image_filename,
                            )

                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=kwargs.get("llm_model"),
                                prompt=kwargs.get("llm_prompt"),
                            )
                        except Exception:
                            pass

                    # Try OCR if LLM failed or not available
                    ocr_text = ""
                    if not llm_description and ocr_service:
                        try:
                            image_stream.seek(0)
                            ocr_result = ocr_service.extract_text(image_stream)
                            if ocr_result.text.strip():
                                ocr_text = ocr_result.text.strip()
                        except Exception:
                            pass

                    # Format extracted content using unified OCR block format
                    content = (llm_description or ocr_text or "").strip()
                    if content:
                        md_content += f"\n*[Image OCR]\n{content}\n[End OCR]*\n"

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\\n"
                    else:
                        md_content += shape.text + "\\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\\n\\n### Notes:\\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        import html

        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\\n\\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\\n\\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            if "unsupported plot type" in str(e):
                return "\\n\\n[unsupported chart]\\n\\n"
        except Exception:
            return "\\n\\n[unsupported chart]\\n\\n"
